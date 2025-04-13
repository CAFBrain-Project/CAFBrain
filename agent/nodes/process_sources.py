from agent.state import State
from agent.utils import save_img
from agent.models.ocr import OCR
from agent.models.image_caption import ImageCaption

import io
import fitz
from pptx import Presentation
from PIL import Image as PILImage

from langchain.docstore.document import Document
from langchain.text_splitter import RecursiveCharacterTextSplitter


ocr = OCR()
image_caption = ImageCaption()
text_splitter = RecursiveCharacterTextSplitter(chunk_size = 500, chunk_overlap = 50)


def process_image(image_path: str, file_type = "image", **kwargs):
    documents = []
    
    # OCR
    ocr_result = ocr.forward(image_path)
    documents.extend([
        Document(
            page_content = chunk, 
            metadata = {
                "file_name": image_path.split("/")[-1], 
                "file_type": file_type, 
                "extraction_type": "OCR",
                **kwargs
            }
        ) for chunk in text_splitter.split_text(ocr_result)
    ])

    # Image Captioning
    image_caption_result = image_caption.forward(image_path)
    documents.extend([
        Document(
            page_content = chunk, 
            metadata = {
                "file_name": image_path.split("/")[-1],
                "file_type": file_type, 
                "extraction_type": "Image Caption",
                **kwargs
            }
        ) for chunk in text_splitter.split_text(image_caption_result)
    ])

    text = "**OCR**: " + ocr_result + "\n\n**Image Caption**: " + image_caption_result + "\n"
    return text, documents

def process_pdf(pdf_path: str):
    doc = fitz.open(pdf_path)

    extracted_text = ""
    documents = []
    
    for page_num, page in enumerate(doc):
        text = page.get_text("text") 
        image_list = page.get_images(full = True) # TODO: Read XML based images too

        documents.extend([
            Document(
                page_content = chunk,
                metadata = {
                    "file_name": pdf_path.split("/")[-1],
                    "page_num": page_num,
                    "file_type": "PDF",
                    "extraction_type": "Text"
                } 
            ) for chunk in text_splitter.split_text(text)
        ])

        images = []
        for img_index, img in enumerate(image_list):
            xref = img[0]  # Xref is the reference to the image
            base_image = doc.extract_image(xref)
            img_bytes = base_image["image"]  # Image as byte data
            img_ext = base_image["ext"]
            
            img = PILImage.open(io.BytesIO(img_bytes))
            image_path = save_img(img, img_ext)

            txt, docs = process_image(image_path, file_type = "PDF", page_num = page_num)
            text += "\n" + txt
            documents.extend(docs)

        extracted_text += text + "\n"
    
    return extracted_text, documents

def process_ppt(ppt_path: str):
    extracted_text = ""
    documents = []
    
    presentation = Presentation(ppt_path)
    
    for slide_num, slide in enumerate(presentation.slides):
        text = []
        images = []
        
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
                
            elif hasattr(shape, "image"):
                image = shape.image
                img_bytes = image.blob
                img_ext = image.ext
                images.append(img_bytes)

                img = PILImage.open(io.BytesIO(img_bytes))
                image_path = save_img(img, img_ext)

                txt, _ = process_image(image_path)
                text.append(txt)

        extracted_text += "\n".join(text) + "\n"
        documents.extend([
            Document(
                page_content = chunk, 
                metadata = {
                    "file_name": ppt_path.split("/")[-1],
                    "slide_num": slide_num,
                    "file_type": "PPT",
                    "extraction_type": "Text + Images"
                } 
            ) for chunk in text_splitter.split_text("\n".join(text))
        ])

    return extracted_text, documents

def process_text(text_file_path: str):
    with open(text_file_path, "r") as fp:
        content = [line.strip() for line in fp.read_lines()]

    documents = [
        Document(
            page_content = chunk,
            metadata = {
                "file_name": text_file_path.split("/")[-1],
                "file_type": "Text File",
                "extraction_type": "Text"
            } 
        ) for chunk in text_splitter.split_text(content)
    ]

    return "\n".join(content), documents

def process_sources(state: State):
    extracted_texts = state.get("extracted_texts", [])
    documents = state.get("documents", [])
    prev_doc_count = state.get("prev_doc_count", 0)
    source_file_paths = state.get("source_file_paths", [])

    # print('To check if processing updated:', len(source_file_paths), prev_doc_count)

    for file_path in source_file_paths[prev_doc_count:]:
        file_extension = (file_path.split('.')[-1]).lower()
        if file_extension in ["png", "jpg", "jpeg"]:
            text, docs = process_image(image_path = file_path)

        elif file_extension in ["txt", "docx"]:
            text, docs = process_text(text_path = file_path)
            
        elif file_extension == "pdf":
            text, docs = process_pdf(pdf_path = file_path)
            
        elif file_extension in ["ppt", "pptx"]:
            text, docs = process_ppt(ppt_path = file_path)

        extracted_texts.append(text)
        documents.extend(docs)

    print("Processed sources")
    
    return state | {"extracted_texts": extracted_texts, "documents": documents}