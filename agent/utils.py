import os
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from fpdf import FPDF
from io import BytesIO
from unidecode import unidecode

# Paths
TEMP_DIR = "temp"
LOGO_PATH = "/Users/nishchalmn/Downloads/CAFB/CAFBrain/logo.png"
FONT_DIR = "/Users/nishchalmn/Downloads/CAFB/Assistant/static"

ASSISTANT_FONTS = {
    "Regular": os.path.join(FONT_DIR, "Assistant-Regular.ttf"),
    "SemiBold": os.path.join(FONT_DIR, "Assistant-SemiBold.ttf"),
    "Medium": os.path.join(FONT_DIR, "Assistant-Medium.ttf"),
    "Light": os.path.join(FONT_DIR, "Assistant-Light.ttf"),
    "ExtraLight": os.path.join(FONT_DIR, "Assistant-ExtraLight.ttf"),
    "Bold": os.path.join(FONT_DIR, "Assistant-Bold.ttf"),
    "ExtraBold": os.path.join(FONT_DIR, "Assistant-ExtraBold.ttf")
}

# Colors
PRIMARY_GREEN = RGBColor(101, 141, 27)
DARK_GREEN = RGBColor(21, 71, 52)
GRAY = RGBColor(63, 68, 67)
LIGHT_TAUPE = RGBColor(242, 240, 235)
ORANGE = RGBColor(220, 122, 0)
WHITE = RGBColor(255, 255, 255)
SAGE_GREEN = RGBColor(154, 171, 137)

# Font Sizes
TITLE_FONT_SIZE = Pt(38)
SUBTITLE_FONT_SIZE = Pt(28)
SUBHEAD_FONT_SIZE = Pt(13)
BODY_FONT_SIZE = Pt(24)
PDF_HEADER_FONT_SIZE = 15
PDF_BODY_FONT_SIZE = 12

SECTION_COLORS = {
    "impact": ORANGE,
    "budget": DARK_GREEN,
    "plan": PRIMARY_GREEN,
    "normal": LIGHT_TAUPE,
    "title": SAGE_GREEN
}

def save_img(img, img_ext):
    os.makedirs(TEMP_DIR, exist_ok=True)
    file_name = os.path.join(TEMP_DIR, f"{len(os.listdir(TEMP_DIR))}.{img_ext}")
    img.save(file_name)
    return file_name

def json_to_blog_html(json_input: str, output_file: str = "blog_post.html") -> str:
    data = json.loads(json_input)
    headline = data.get("headline", "")
    introduction = data.get("introduction", "")
    body = data.get("body", [])
    conclusion = data.get("conclusion", "")

    html_content = f"""
    <!DOCTYPE html>
    <html lang=\"en\">
    <head>
        <meta charset=\"UTF-8\">
        <meta name=\"viewport\" content=\"width=device-width, initial-scale=1.0\">
        <title>{headline}</title>
        <style>
            @font-face {{
                font-family: 'Assistant';
                src: url('{ASSISTANT_FONTS['Regular']}');
            }}
            body {{
                font-family: 'Assistant', sans-serif;
                background-color: #f2f0eb;
                color: #3f4443;
                max-width: 800px;
                margin: 2rem auto;
                padding: 1rem;
                line-height: 1.6;
            }}
            h1 {{
                font-weight: 800;
                font-size: 2.25rem;
                color: #658d1b;
            }}
            h2 {{
                font-weight: 700;
                font-size: 1.5rem;
                color: #154734;
                text-transform: uppercase;
            }}
            p {{
                font-size: 1rem;
                margin-bottom: 1rem;
            }}
            .highlight {{
                background-color: #f2f0eb;
                padding: 0.5rem;
                border-left: 5px solid #dc7a00;
            }}
            .logo-header {{
                text-align: center;
                margin-bottom: 2rem;
            }}
        </style>
    </head>
    <body>
        <div class="logo-header">
            <img src="/Users/nishchalmn/Downloads/CAFB/CAFBrain/logo.png" alt="CAFB Logo" width="120">
        </div>
        <h1>{headline}</h1>
        <p><strong>{introduction}</strong></p>
    """
    for paragraph in body:
        html_content += f"<p>{paragraph}</p>\n"
    html_content += f"<p class='highlight'><strong>{conclusion}</strong></p>\n</body></html>"
    return html_content.strip()

def add_logo(slide):
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, Inches(0.1), Inches(0.1), width=Inches(1))

def apply_slide_style(slide, section="normal"):
    if section in ["impact", "budget", "plan"]:
        color = SECTION_COLORS[section]
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color
        for shape in slide.shapes:
            if shape.has_text_frame:
                shape.text_frame.text = shape.text_frame.text
                for p in shape.text_frame.paragraphs:
                    p.font.size = TITLE_FONT_SIZE
                    p.font.bold = True
                    p.font.color.rgb = WHITE

    elif section == 'title':
        color = SECTION_COLORS[section]
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color
        for shape in slide.shapes:
            if shape.has_text_frame:
                shape.text_frame.text = shape.text_frame.text
                for p in shape.text_frame.paragraphs:
                    p.font.size = TITLE_FONT_SIZE
                    p.font.bold = True
                    p.font.color.rgb = WHITE
    else:
        color = SECTION_COLORS[section]
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color


def json_to_presentation(json_input: str, output_file: str = "content.pptx") -> BytesIO:
    data = json.loads(json_input)
    prs = Presentation()
    TITLE_SLIDE_LAYOUT = 0
    TITLE_AND_CONTENT_LAYOUT = 1

    for slide_data in data:
        is_title_slide = "subtitle" in slide_data
        layout = TITLE_SLIDE_LAYOUT if is_title_slide else TITLE_AND_CONTENT_LAYOUT
        slide = prs.slides.add_slide(prs.slide_layouts[layout])
        slide.shapes.title.text = slide_data["title"]

        if is_title_slide:
            slide.placeholders[1].text = slide_data["subtitle"]
            apply_slide_style(slide, "title")
        else:
            bullets = slide_data.get("bullets", [])
            section = "normal"
            if "impact" in slide_data["title"].lower():
                section = "impact"
            elif "budget" in slide_data["title"].lower():
                section = "budget"
            elif "plan" in slide_data["title"].lower() or "implementation" in slide_data["title"].lower():
                section = "plan"
            content = slide.placeholders[1].text_frame
            content.clear()
            for bullet in bullets:
                p = content.add_paragraph()
                p.text = bullet
                p.level = 0
            apply_slide_style(slide, section)

        add_logo(slide)

    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer


class ProposalPDF(FPDF):
    def header(self):
        if os.path.exists(LOGO_PATH):
            self.image(LOGO_PATH, 10, 8, 33)
        self.set_xy(50, 10)
        self.set_font("Assistant", "B", PDF_HEADER_FONT_SIZE)
        self.set_text_color(*DARK_GREEN)
        self.cell(0, 10, "Grant Proposal", ln=True, align="R")
        self.ln(10)

    def chapter_title(self, title):
        self.set_font("Assistant", "B", 14)
        self.set_text_color(*PRIMARY_GREEN)
        self.multi_cell(0, 10, title)
        self.ln(2)

    def chapter_body(self, text):
        self.set_font("Assistant", "", PDF_BODY_FONT_SIZE)
        self.set_text_color(*GRAY)
        self.multi_cell(0, 8, text)
        self.ln(4)

def json_to_pdf(json_input: str) -> BytesIO:
    data = json.loads(json_input)
    pdf = ProposalPDF()
    for style, path in ASSISTANT_FONTS.items():
        font_weight = "" if style in ["Regular", "Light", "ExtraLight"] else "B"
        pdf.add_font("Assistant", font_weight, path, uni=True)
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()

    for section in data:
        heading = unidecode(section.get("heading", ""))
        content = unidecode(section.get("content", ""))
        pdf.chapter_title(heading)
        pdf.chapter_body(content)

    pdf_bytes = pdf.output(dest='S')
    buffer = BytesIO()
    buffer.write(pdf_bytes)
    buffer.seek(0)
    return buffer
