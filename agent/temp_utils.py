# import os
# import json
# from pptx import Presentation
# from pptx.util import Inches
# from io import BytesIO
# from fpdf import FPDF
# from unidecode import unidecode
# from pptx.util import Pt

# def save_img(img, img_ext):
#     if not os.path.exists("temp"):
#         os.mkdir("temp")
#     file_name = os.path.join("temp", str(len(os.listdir("temp"))) + "." + img_ext)
#     img.save(file_name)
#     return file_name

# def generate_content(template_format: str, data: dict) -> dict:
#     """
#     Generates content based on the specified template format.

#     Args:
#         template_format (str): The format to generate content for ('blog', 'ppt', 'pdf', 'grant').
#         data (dict): Input data for content generation.

#     Returns:
#         dict: Generated content in a structured format.  Returns an empty dict if the template_format is not supported.
#     """
#     if template_format == "blog":
#         return {
#             "headline": data.get("headline", "Blog Post"),
#             "introduction": data.get("introduction", ""),
#             "body": data.get("body", []),
#             "conclusion": data.get("conclusion", "")
#         }
#     elif template_format == "ppt":
#         slides = []
#         for slide_data in data.get("slides", []):
#             slide = {
#                 "title": slide_data.get("title", ""),
#                 "subtitle": slide_data.get("subtitle", ""),
#                 "bullets": slide_data.get("bullets", []),
#                 "images": slide_data.get("images", []) # Added image support
#             }
#             slides.append(slide)
#         return {"slides": slides}

#     elif template_format == "pdf":
#         sections = []
#         for section_data in data.get("sections", []):
#             section = {
#                 "heading": section_data.get("heading", ""),
#                 "content": section_data.get("content", "")
#             }
#             sections.append(section)
#         return {"sections": sections}
#     elif template_format == "grant":
#         sections = []
#         for section_data in data.get("sections", []):
#             section = {
#                 "heading": section_data.get("heading", ""),
#                 "content": section_data.get("content", ""),
#                 "budget": section_data.get("budget", None)  # Example: Include budget info
#             }
#             sections.append(section)
#         return {"sections": sections}
#     else:
#         return {}

# def json_to_blog_html(json_input: dict, output_file: str = "blog_post.html") -> str:
#     """
#     Converts JSON data to a blog post HTML format.

#     Args:
#         json_input (dict):  Dictionary, not a JSON string.
#         output_file (str): The name of the output HTML file.

#     Returns:
#         str: The generated HTML content.
#     """
#     # Extract content
#     headline = json_input.get("headline", "")
#     introduction = json_input.get("introduction", "")
#     body = json_input.get("body", [])
#     conclusion = json_input.get("conclusion", "")

#     # Build HTML content
#     html_content = f"""
#     <!DOCTYPE html>
#     <html lang="en">
#     <head>
#         <meta charset="UTF-8">
#         <meta name="viewport" content="width=device-width, initial-scale=1.0">
#         <title>{headline}</title>
#         <style>
#             body {{
#                 font-family: Arial, sans-serif;
#                 line-height: 1.6;
#                 max-width: 800px;
#                 margin: 2rem auto;
#                 padding: 1rem;
#                 background-color: #f9f9f9;
#                 color: #333;
#             }}
#             h1 {{
#                 color: #004080;
#             }}
#             p {{
#                 margin-bottom: 1rem;
#             }}
#         </style>
#     </head>
#     <body>
#         <h1>{headline}</h1>
#         <p>{introduction}</p>
#     """

#     for paragraph in body:
#         html_content += f"<p>{paragraph}</p>\n"

#     html_content += f"<p><strong>{conclusion}</strong></p>\n"

#     html_content += """
#     </body>
#     </html>
#     """

#     return html_content.strip()


# def json_to_presentation(json_input: dict, template_file: str = None) -> BytesIO:
#     """
#     Converts JSON data to a PowerPoint presentation, optionally using a template.

#     Args:
#         json_input (dict):  Dictionary, not a JSON string.
#         template_file (str, optional): Path to a PowerPoint template file.  Defaults to None.

#     Returns:
#         BytesIO: A byte stream containing the PowerPoint presentation.
#     """
#     if template_file:
#         prs = Presentation(template_file)  # Load the template
#     else:
#         prs = Presentation()  # Create a new presentation

#     # Define layout indexes.  These might need to be adjusted based on the template.
#     #  It's CRITICAL to inspect your template_file and determine the correct layout IDs.
#     #  The code assumes default layouts if no template is provided.  With a template,
#     #  you'll likely need to change these.
#     TITLE_SLIDE_LAYOUT = 0
#     TITLE_AND_CONTENT_LAYOUT = 1
#     BLANK_SLIDE_LAYOUT = 6 # Added for Blank slide.

#     for slide_data in json_input.get("slides", []): # Access the slides from the dict
#         layout_type = slide_data.get("layout", "title_and_content") # default
#         if layout_type == "title_slide":
#             slide_layout = prs.slide_layouts[TITLE_SLIDE_LAYOUT]
#         elif layout_type == "title_and_content":
#              slide_layout = prs.slide_layouts[TITLE_AND_CONTENT_LAYOUT]
#         elif layout_type == "blank":
#             slide_layout = prs.slide_layouts[BLANK_SLIDE_LAYOUT] # Use blank layout
#         else:
#             slide_layout = prs.slide_layouts[TITLE_AND_CONTENT_LAYOUT] # Fallback

#         slide = prs.slides.add_slide(slide_layout)

#         if layout_type == "title_slide":
#             if "title" in slide_data:
#                 slide.shapes.title.text = slide_data["title"]
#             if "subtitle" in slide_data:
#                 try:
#                     slide.placeholders[1].text = slide_data["subtitle"]
#                 except IndexError:
#                     print(f"Warning: Subtitle placeholder not found on title slide. Slide data: {slide_data}")

#         elif layout_type == "title_and_content":
#             if "title" in slide_data:
#                 slide.shapes.title.text = slide_data["title"]
#             if "bullets" in slide_data:
#                 content = slide.placeholders[1].text_frame
#                 content.clear()
#                 for bullet in slide_data["bullets"]:
#                     p = content.add_paragraph()
#                     p.text = bullet
#                     p.level = 0
#             if "images" in slide_data: #add images
#                 for img_data in slide_data["images"]:
#                     img_path = img_data["path"]
#                     left = Inches(img_data.get("left", 1))
#                     top = Inches(img_data.get("top", 1))
#                     width = Inches(img_data.get("width", 3))
#                     height = Inches(img_data.get("height", 2))
#                     try:
#                         slide.shapes.add_picture(img_path, left, top, width, height)
#                     except Exception as e:
#                         print(f"Error adding image {img_path}: {e}.  Slide data: {slide_data}")
#         elif layout_type == "blank":
#              if "images" in slide_data: #add images
#                 for img_data in slide_data["images"]:
#                     img_path = img_data["path"]
#                     left = Inches(img_data.get("left", 1))
#                     top = Inches(img_data.get("top", 1))
#                     width = Inches(img_data.get("width", 3))
#                     height = Inches(img_data.get("height", 2))
#                     try:
#                         slide.shapes.add_picture(img_path, left, top, width, height)
#                     except Exception as e:
#                         print(f"Error adding image {img_path}: {e}.  Slide data: {slide_data}")

#     # Save to in-memory buffer
#     buffer = BytesIO()
#     prs.save(buffer)
#     buffer.seek(0)
#     return buffer


# class ProposalPDF(FPDF):
#     """
#     Custom PDF class for generating grant proposals.  This class can be extended
#     to further customize the PDF's appearance.  For example, you could add a logo,
#     footer, or custom fonts.
#     """
#     def __init__(self, **kwargs):
#         super().__init__(**kwargs)
#         self.title = "Grant Proposal"  # Default title. Can be overridden.

#     def header(self):
#         """Adds a header to the PDF."""
#         self.set_font("Helvetica", "B", 12)
#         self.cell(0, 10, self.title, ln=True, align="C")  # Use self.title
#         self.ln(5)

#     def chapter_title(self, title):
#         """Adds a chapter title to the PDF."""
#         self.set_font("Helvetica", "B", 14)
#         self.set_text_color(0)
#         self.multi_cell(0, 10, title)
#         self.ln(2)

#     def chapter_body(self, text):
#         """Adds the body text of a chapter to the PDF."""
#         self.set_font("Helvetica", "", 12)
#         self.set_text_color(50)
#         self.multi_cell(0, 8, text)
#         self.ln(4)

#     def add_budget(self, budget_data):
#         """Adds a budget section to the PDF.

#         Args:
#             budget_data (dict):  A dictionary containing budget information.
#                 Example:
#                 {
#                     "items": [
#                         {"item": "Personnel", "amount": 50000},
#                         {"item": "Materials", "amount": 10000},
#                         {"item": "Travel", "amount": 5000}
#                     ],
#                     "total": 65000
#                 }
#         """
#         self.set_font("Helvetica", "B", 14)
#         self.chapter_title("Budget")
#         self.set_font("Helvetica", "", 12)

#         if not budget_data:
#             self.chapter_body("No budget information provided.")
#             return

#         # Create a table-like structure.
#         line_height = 6
#         col_width = [80, 30]  # Adjust column widths as needed.

#         # Header
#         self.cell(col_width[0], line_height, "Item", border=1)
#         self.cell(col_width[1], line_height, "Amount", border=1)
#         self.ln()

#         # Data rows
#         total = 0
#         if "items" in budget_data:
#             for item in budget_data["items"]:
#                 self.cell(col_width[0], line_height, item.get("item", ""), border=1)
#                 amount = item.get("amount", 0)
#                 self.cell(col_width[1], line_height, str(amount), border=1)
#                 self.ln()
#                 total += amount
#         else:
#              self.cell(col_width[0], line_height, "Total", border=1)
#              total = budget_data.get("total", 0)
#              self.cell(col_width[1], line_height, str(total), border=1)
#              self.ln()
#         # Total
#         self.set_font("Helvetica", "B", 12)
#         self.cell(col_width[0], line_height, "Total", border=1)
#         self.cell(col_width[1], line_height, str(total), border=1)
#         self.ln()
#         self.ln(4)



# def json_to_pdf(json_input: dict, template_file: str = None) -> BytesIO:
#     """
#     Converts JSON data to a PDF, optionally using a template (for grant proposals).

#     Args:
#         json_input (dict):  Dictionary, not a JSON string.
#         template_file (str, optional): Path to a PDF template file (not fully supported,
#             but can influence behavior). Defaults to None.
#     Returns:
#         BytesIO: A byte stream containing the PDF.
#     """
#     if template_file:
#         #  PDF templates are harder to use directly.  You might use it to get
#         #  styles, or page dimensions, but not usually the full layout.
#         #  For now, we'll just note that a template was provided.
#         print(f"Using template file: {template_file} (Note: Full template support is limited.)")

#     pdf = ProposalPDF()  # Use our custom class
#     pdf.set_auto_page_break(auto=True, margin=15)
#     pdf.add_page()

#     # Set document title.
#     pdf.title = json_input.get("title", "Grant Proposal") # set the title
#     pdf.set_title(pdf.title)

#     for section in json_input.get("sections", []): # get sections from json
#         heading = unidecode(section.get("heading", ""))
#         content = unidecode(section.get("content", ""))
#         pdf.chapter_title(heading)
#         pdf.chapter_body(content)
#         if "budget" in section:
#             pdf.add_budget(section["budget"]) # Add budget if present

#     pdf_bytes = pdf.output(dest='S').encode('latin1')  # Get PDF as bytes
#     buffer = BytesIO()
#     buffer.write(pdf_bytes)
#     buffer.seek(0)
#     return buffer

# def convert_content(input_json: str, output_format: str, template_file: str = None) -> BytesIO | str:
#     """
#     Converts JSON data to the specified output format, optionally using a template.

#     Args:
#         input_json (str): A JSON string containing the data to convert.
#         output_format (str): The desired output format ('blog', 'ppt', 'pdf', 'grant').
#         template_file (str, optional): Path to a template file (for PPT or PDF).
#             Defaults to None.

#     Returns:
#         BytesIO | str: The converted content as a byte stream (for PPT, PDF) or a string (for blog),
#                        or None on error.
#     """
#     try:
#         data = json.loads(input_json)
#     except json.JSONDecodeError:
#         print("Error: Invalid JSON input.")
#         return None

#     # Use the generate_content function
#     content_data = generate_content(output_format, data)

#     if not content_data:
#         print(f"Error: Unsupported output format: {output_format}")
#         return None

#     if output_format == "blog":
#         return json_to_blog_html(content_data)  # Pass the dictionary, not the string
#     elif output_format == "ppt":
#         return json_to_presentation(content_data, template_file)
#     elif output_format == "pdf" or output_format == "grant":
#         return json_to_pdf(content_data, template_file)
#     else:
#         print(f"Error: Unsupported output format: {output_format}")
#         return None

# def main():
#     """
#     Main function to demonstrate the conversion process.
#     """
#     # Example JSON data (simulated)
#     example_data = {
#         "blog": {
#             "headline": "My Awesome Blog Post",
#             "introduction": "This is an introduction to my blog post.",
#             "body": [
#                 "This is the first paragraph of the body.",
#                 "This is the second paragraph of the body.",
#                 "Here's another paragraph with some <em>emphasis</em>.",
#             ],
#             "conclusion": "In conclusion, this blog post is awesome!"
#         },
#         "ppt": {
#             "slides": [
#                 {
#                     "layout": "title_slide",
#                     "title": "Presentation Title",
#                     "subtitle": "Subtitle of the presentation",
#                 },
#                 {
#                     "layout": "title_and_content",
#                     "title": "Slide 1",
#                     "bullets": ["Bullet 1", "Bullet 2", "Bullet 3"],
#                 },
#                 {
#                     "layout": "title_and_content",
#                     "title": "Slide 2",
#                     "bullets": ["Point A", "Point B", "Point C"],
#                     "images": [
#                         {"path": "temp/image1.jpg", "left": 1, "top": 1, "width": 3, "height": 2},
#                         {"path": "temp/image2.png", "left": 4, "top": 3, "width": 2, "height": 2}
#                     ]
#                 },
#                 {
#                     "layout": "blank",
#                     "images": [
#                         {"path": "temp/full_page_image.jpg", "left": 0, "top": 0, "width": 10, "height": 7.5}
#                     ]
#                 }
#             ]
#         },
#         "pdf": {
#             "title": "My Report",  # Added title for PDF
#             "sections": [
#                 {
#                     "heading": "Section 1: Introduction",
#                     "content": "This is the introduction section of the report.  It provides background information."
#                 },
#                 {
#                     "heading": "Section 2: Methods",
#                     "content": "This section describes the methods used in the study."
#                 },
#                 {
#                   "heading": "Section 3: Results",
#                   "content": "This section shows the results of the study.",
#                 }
#             ]
#         },
#         "grant": {
#             "title": "Research Grant Proposal",
#             "sections": [
#                 {
#                     "heading": "Project Summary",
#                     "content": "This project aims to...",
#                 },
#                 {
#                     "heading": "Project Description",
#                     "content": "Detailed description of the research...",
#                 },
#                 {
#                     "heading": "Budget",
#                     "content": "Details of the budget.",
#                     "budget": {
#                         "items": [
#                             {"item": "Personnel", "amount": 50000},
#                             {"item": "Materials", "amount": 10000},
#                             {"item": "Travel", "amount": 5000}
#                         ],
#                         "total": 65000
#                     }
#                 },
#                 {
#                     "heading": "Evaluation Plan",
#                     "content": "How the project will be evaluated."
#                 }
#             ]
#         }
#     }

#     # Create dummy images for PPT testing
#     from PIL import Image
#     img1 = Image.new('RGB', (100, 100), color='red')
#     img2 = Image.new('RGB', (100, 100), color='green')
#     img3 = Image.new('RGB', (800, 600), color='blue') # full page image
#     save_img(img1, 'jpg')
#     save_img(img2, 'png')
#     save_img(img3, 'jpg')


#     # Example usage:
#     blog_html = convert_content(json.dumps(example_data["blog"]), "blog")
#     if blog_html:
#         with open("blog_output.html", "w") as f:
#             f.write(blog_html)
#         print("Blog HTML generated successfully (blog_output.html).")

#     ppt_bytes = convert_content(json.dumps(example_data["ppt"]), "ppt", template_file="CAFB_SlideDeck_Template.pptx") # Pass the template
#     if ppt_bytes:
#         with open("presentation_output.pptx", "wb") as f:
#             f.write(ppt_bytes.getvalue())
#         print("PowerPoint presentation generated successfully (presentation_output.pptx).")

#     pdf_bytes = convert_content(json.dumps(example_data["pdf"]), "pdf")
#     if pdf_bytes:
#         with open("report_output.pdf", "wb") as f:
#             f.write(pdf_bytes.getvalue())
#         print("PDF report generated successfully (report_output.pdf).")

#     grant_pdf_bytes = convert_content(json.dumps(example_data["grant"]), "grant", template_file="grant_proposal_template.pdf")
#     if grant_pdf_bytes:
#         with open("grant_proposal_output.pdf", "wb") as f:
#             f.write(grant_pdf_bytes.getvalue())
#         print("Grant proposal PDF generated successfully (grant_proposal_output.pdf).")



# if __name__ == "__main__":
#     main()


# import os
# import json
# from pptx import Presentation
# from fpdf import FPDF
# from pptx.util import Inches
# from io import BytesIO
# from unidecode import unidecode


# def save_img(img, img_ext):
#     if not os.path.exists("temp"):
#         os.mkdir("temp")

#     file_name = os.path.join("temp", str(len(os.listdir("temp"))) + "." + img_ext)
#     img.save(file_name)

#     return file_name

# def json_to_blog_html(json_input: str, output_file: str = "blog_post.html") -> str:
#     data = json.loads(json_input)
    
#     # Extract content
#     headline = data.get("headline", "")
#     introduction = data.get("introduction", "")
#     body = data.get("body", [])
#     conclusion = data.get("conclusion", "")

#     # Build HTML content
#     html_content = f"""
#     <!DOCTYPE html>
#     <html lang="en">
#     <head>
#         <meta charset="UTF-8">
#         <meta name="viewport" content="width=device-width, initial-scale=1.0">
#         <title>{headline}</title>
#         <style>
#             body {{
#                 font-family: Arial, sans-serif;
#                 line-height: 1.6;
#                 max-width: 800px;
#                 margin: 2rem auto;
#                 padding: 1rem;
#                 background-color: #f9f9f9;
#                 color: #333;
#             }}
#             h1 {{
#                 color: #004080;
#             }}
#             p {{
#                 margin-bottom: 1rem;
#             }}
#         </style>
#     </head>
#     <body>
#         <h1>{headline}</h1>
#         <p>{introduction}</p>
#     """

#     for paragraph in body:
#         html_content += f"<p>{paragraph}</p>\n"

#     html_content += f"<p><strong>{conclusion}</strong></p>\n"

#     html_content += """
#     </body>
#     </html>
#     """

#     # Save to file
#     # with open(output_file, "w", encoding="utf-8") as f:
#     #     f.write(html_content.strip())

#     return html_content.strip()



# def json_to_presentation(json_input: str, output_file: str = "content.pptx") -> BytesIO:

#     json_input = json.loads(json_input)
#     prs = Presentation()

#     # Define layout indexes
#     TITLE_SLIDE_LAYOUT = 0
#     TITLE_AND_CONTENT_LAYOUT = 1

#     for slide_data in json_input:
#         if "subtitle" in slide_data:
#             slide_layout = prs.slide_layouts[TITLE_SLIDE_LAYOUT]
#             slide = prs.slides.add_slide(slide_layout)
#             slide.shapes.title.text = slide_data["title"]
#             slide.placeholders[1].text = slide_data["subtitle"]
#         else:
#             slide_layout = prs.slide_layouts[TITLE_AND_CONTENT_LAYOUT]
#             slide = prs.slides.add_slide(slide_layout)
#             slide.shapes.title.text = slide_data["title"]
#             content = slide.placeholders[1].text_frame
#             content.clear()

#             for bullet in slide_data.get("bullets", []):
#                 p = content.add_paragraph()
#                 p.text = bullet
#                 p.level = 0

#     # Save to in-memory buffer
#     buffer = BytesIO()
#     prs.save(buffer)
#     buffer.seek(0)
#     return buffer

# class ProposalPDF(FPDF):
#     def header(self):
#         self.set_font("Helvetica", "B", 12)
#         self.cell(0, 10, "Grant Proposal", ln=True, align="C")
#         self.ln(5)

#     def chapter_title(self, title):
#         self.set_font("Helvetica", "B", 14)
#         self.set_text_color(0)
#         self.multi_cell(0, 10, title)
#         self.ln(2)

#     def chapter_body(self, text):
#         self.set_font("Helvetica", "", 12)
#         self.set_text_color(50)
#         self.multi_cell(0, 8, text)
#         self.ln(4)

# def json_to_pdf(json_input: str) -> BytesIO:
#     data = json.loads(json_input)
#     pdf = ProposalPDF()
#     pdf.set_auto_page_break(auto=True, margin=15)
#     pdf.add_page()

#     for section in data:
#         heading = unidecode(section.get("heading", ""))
#         content = unidecode(section.get("content", ""))
#         pdf.chapter_title(heading)
#         pdf.chapter_body(content)

#     pdf_bytes = pdf.output(dest='S').encode('latin1')  # Get PDF as bytes
#     buffer = BytesIO()
#     buffer.write(pdf_bytes)
#     buffer.seek(0)
#     return buffer





import os
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from fpdf import FPDF
from io import BytesIO
from unidecode import unidecode

# Paths
TEMP_DIR = "temp"
LOGO_PATH = "assets/cafb_logo.png"
FONT_DIR = "fonts"
ASSISTANT_FONTS = {
    "Regular": os.path.join(FONT_DIR, "Assistant-Regular.ttf"),
    "SemiBold": os.path.join(FONT_DIR, "Assistant-SemiBold.ttf"),
    "Medium": os.path.join(FONT_DIR, "Assistant-Medium.ttf"),
    "Light": os.path.join(FONT_DIR, "Assistant-Light.ttf"),
    "ExtraLight": os.path.join(FONT_DIR, "Assistant-ExtraLight.ttf"),
    "Bold": os.path.join(FONT_DIR, "Assistant-Bold.ttf"),
    "ExtraBold": os.path.join(FONT_DIR, "Assistant-ExtraBold.ttf")
}

# Colors
PRIMARY_GREEN = RGBColor(101, 141, 27)
DARK_GREEN = RGBColor(21, 71, 52)
GRAY = RGBColor(63, 68, 67)
LIGHT_TAUPE = RGBColor(242, 240, 235)
ORANGE = RGBColor(220, 122, 0)
WHITE = RGBColor(255, 255, 255)

# Font Sizes
TITLE_FONT_SIZE = Pt(38)
SUBHEAD_FONT_SIZE = Pt(13)
BODY_FONT_SIZE = Pt(24)
PDF_HEADER_FONT_SIZE = 15
PDF_BODY_FONT_SIZE = 12

SECTION_COLORS = {
    "impact": ORANGE,
    "budget": DARK_GREEN,
    "plan": PRIMARY_GREEN
}

def save_img(img, img_ext):
    os.makedirs(TEMP_DIR, exist_ok=True)
    file_name = os.path.join(TEMP_DIR, f"{len(os.listdir(TEMP_DIR))}.{img_ext}")
    img.save(file_name)
    return file_name

def json_to_blog_html(json_input: str, output_file: str = "blog_post.html") -> str:
    data = json.loads(json_input)
    headline = data.get("headline", "")
    introduction = data.get("introduction", "")
    body = data.get("body", [])
    conclusion = data.get("conclusion", "")

    html_content = f"""
    <!DOCTYPE html>
    <html lang=\"en\">
    <head>
        <meta charset=\"UTF-8\">
        <meta name=\"viewport\" content=\"width=device-width, initial-scale=1.0\">
        <title>{headline}</title>
        <style>
            @font-face {{
                font-family: 'Assistant';
                src: url('{ASSISTANT_FONTS['Regular']}');
            }}
            body {{
                font-family: 'Assistant', sans-serif;
                background-color: #f2f0eb;
                color: #3f4443;
                max-width: 800px;
                margin: 2rem auto;
                padding: 1rem;
                line-height: 1.6;
            }}
            h1 {{
                font-weight: 800;
                font-size: 2.25rem;
                color: #658d1b;
            }}
            h2 {{
                font-weight: 700;
                font-size: 1.5rem;
                color: #154734;
                text-transform: uppercase;
            }}
            p {{
                font-size: 1rem;
                margin-bottom: 1rem;
            }}
            .highlight {{
                background-color: #f2f0eb;
                padding: 0.5rem;
                border-left: 5px solid #dc7a00;
            }}
            .logo-header {{
                text-align: center;
                margin-bottom: 2rem;
            }}
        </style>
    </head>
    <body>
        <div class="logo-header">
            <img src="assets/cafb_logo.png" alt="CAFB Logo" width="120">
        </div>
        <h1>{headline}</h1>
        <p><strong>{introduction}</strong></p>
    """
    for paragraph in body:
        html_content += f"<p>{paragraph}</p>\n"
    html_content += f"<p class='highlight'><strong>{conclusion}</strong></p>\n</body></html>"
    return html_content.strip()

def add_logo(slide):
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, Inches(0.1), Inches(0.1), width=Inches(1))

def apply_slide_style(slide, section=""):
    if section in SECTION_COLORS:
        color = SECTION_COLORS[section]
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color
        for shape in slide.shapes:
            if shape.has_text_frame:
                shape.text_frame.text = shape.text_frame.text
                for p in shape.text_frame.paragraphs:
                    p.font.size = Pt(50)
                    p.font.bold = True
                    p.font.color.rgb = WHITE

def json_to_presentation(json_input: str, output_file: str = "content.pptx") -> BytesIO:
    data = json.loads(json_input)
    prs = Presentation()
    TITLE_SLIDE_LAYOUT = 0
    TITLE_AND_CONTENT_LAYOUT = 1

    for slide_data in data:
        is_title_slide = "subtitle" in slide_data
        layout = TITLE_SLIDE_LAYOUT if is_title_slide else TITLE_AND_CONTENT_LAYOUT
        slide = prs.slides.add_slide(prs.slide_layouts[layout])
        slide.shapes.title.text = slide_data["title"]

        if is_title_slide:
            slide.placeholders[1].text = slide_data["subtitle"]
            apply_slide_style(slide, "impact")
        else:
            bullets = slide_data.get("bullets", [])
            section = ""
            if "impact" in slide_data["title"].lower():
                section = "impact"
            elif "budget" in slide_data["title"].lower():
                section = "budget"
            elif "plan" in slide_data["title"].lower() or "implementation" in slide_data["title"].lower():
                section = "plan"
            content = slide.placeholders[1].text_frame
            content.clear()
            for bullet in bullets:
                p = content.add_paragraph()
                p.text = bullet
                p.level = 0
            apply_slide_style(slide, section)

        add_logo(slide)

    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer

class ProposalPDF(FPDF):
    def header(self):
        if os.path.exists(LOGO_PATH):
            self.image(LOGO_PATH, 10, 8, 33)
        self.set_xy(50, 10)
        self.set_font("Assistant", "B", PDF_HEADER_FONT_SIZE)
        self.set_text_color(*DARK_GREEN)
        self.cell(0, 10, "Grant Proposal", ln=True, align="R")
        self.ln(10)

    def chapter_title(self, title):
        self.set_font("Assistant", "B", 14)
        self.set_text_color(*PRIMARY_GREEN)
        self.multi_cell(0, 10, title)
        self.ln(2)

    def chapter_body(self, text):
        self.set_font("Assistant", "", PDF_BODY_FONT_SIZE)
        self.set_text_color(*GRAY)
        self.multi_cell(0, 8, text)
        self.ln(4)

def json_to_pdf(json_input: str) -> BytesIO:
    data = json.loads(json_input)
    pdf = ProposalPDF()
    for style, path in ASSISTANT_FONTS.items():
        font_weight = "" if style in ["Regular", "Light", "ExtraLight"] else "B"
        pdf.add_font("Assistant", font_weight, path, uni=True)
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()

    for section in data:
        heading = unidecode(section.get("heading", ""))
        content = unidecode(section.get("content", ""))
        pdf.chapter_title(heading)
        pdf.chapter_body(content)

    pdf_bytes = pdf.output(dest='S').encode('latin1')
    buffer = BytesIO()
    buffer.write(pdf_bytes)
    buffer.seek(0)
    return buffer
